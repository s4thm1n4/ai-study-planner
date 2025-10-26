import fitz  # PyMuPDF
import pptx
import pytesseract
from PIL import Image
import io
from fastapi import UploadFile
from typing import Optional # This one is fine

# Note: You must install the Tesseract-OCR application on your system
# for pytesseract to work. https://github.com/tesseract-ocr/tesseract

async def extract_text_from_pdf(contents: bytes) -> str:
    """Extracts text from a PDF file."""
    text = ""
    with fitz.open(stream=contents, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

async def extract_text_from_pptx(contents: bytes) -> str:
    """Extracts text from a PowerPoint file."""
    text = ""
    with io.BytesIO(contents) as f:
        prs = pptx.Presentation(f)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

async def extract_text_from_image(contents: bytes) -> str:
    """Extracts text from an image file using OCR."""
    with io.BytesIO(contents) as f:
        img = Image.open(f)
        text = pytesseract.image_to_string(img)
    return text

async def parse_document(file: UploadFile) -> str:
    """
    Parses an uploaded file (PDF, PPTX, or Image) and extracts all text.
    """
    contents = await file.read()
    content_type = file.content_type

    print(f"[DEBUG] Parsing file: {file.filename}, Type: {content_type}")

    if content_type == "application/pdf":
        return await extract_text_from_pdf(contents)
    elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return await extract_text_from_pptx(contents)
    elif content_type in ["image/png", "image/jpeg", "image/gif"]:
        return await extract_text_from_image(contents)
    elif content_type == "text/plain":
        return contents.decode("utf-8")
    else:
        raise ValueError(f"Unsupported file type: {content_type}")